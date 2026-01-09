# python-etl/pipelines/ppt_pipeline.py

import io
import logging
import hashlib

from typing import List, Dict, Any, Optional
from zipfile import ZipFile

from minio import Minio
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import psycopg2

logger = logging.getLogger(__name__)


def calculate_file_hash(data: bytes) -> str:
    """Return SHA-256 hash of a bytes buffer."""
    return hashlib.sha256(data).hexdigest()


def is_duplicate(pg_conn, content_hash: str) -> bool:
    """Check if this hash already exists in the catalog."""
    if pg_conn is None:
        return False

    cursor = pg_conn.cursor()
    try:
        cursor.execute(
            """
            SELECT 1 FROM minio_data_catalog
            WHERE content_hash = %s
            LIMIT 1
            """,
            (content_hash,),
        )
        return cursor.fetchone() is not None
    except Exception as e:
        logger.warning(f"[ppt] Error checking duplicate: {e}")
        return False
    finally:
        cursor.close()


def _ensure_unstructured_table(pg_conn):
    """Ensure unstructured_documents table exists."""
    cursor = pg_conn.cursor()
    try:
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS unstructured_documents (
                id SERIAL PRIMARY KEY,
                object_name TEXT NOT NULL,
                file_type TEXT,
                text_content TEXT,
                content_hash TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        
        # Try to add UNIQUE constraint if it doesn't exist
        try:
            cursor.execute(
                """
                ALTER TABLE unstructured_documents 
                ADD CONSTRAINT unstructured_documents_content_hash_key 
                UNIQUE (content_hash)
                """
            )
            logger.info("[ppt] Added UNIQUE constraint to content_hash")
        except psycopg2.errors.DuplicateTable:
            pg_conn.rollback()
            logger.debug("[ppt] UNIQUE constraint already exists on content_hash")
        except Exception as e:
            pg_conn.rollback()
            logger.warning(f"[ppt] Could not add UNIQUE constraint: {e}")
        
        pg_conn.commit()
        
    except Exception as e:
        pg_conn.rollback()
        logger.exception(f"[ppt] Failed ensuring unstructured_documents table: {e}")
        raise
    finally:
        cursor.close()


def _save_unstructured_doc(
    pg_conn,
    object_name: str,
    file_type: str,
    text: str,
    content_hash: str,
):
    """Save extracted text to unstructured_documents table."""
    if not pg_conn:
        logger.warning("[ppt] No database connection, skipping unstructured doc save")
        return
    
    _ensure_unstructured_table(pg_conn)
    cursor = pg_conn.cursor()
    
    try:
        # Check if this hash already exists
        cursor.execute(
            """
            SELECT id FROM unstructured_documents 
            WHERE content_hash = %s 
            LIMIT 1
            """,
            (content_hash,)
        )
        existing = cursor.fetchone()
        
        if existing:
            # Update existing record
            cursor.execute(
                """
                UPDATE unstructured_documents
                SET text_content = %s,
                    object_name = %s,
                    file_type = %s,
                    created_at = CURRENT_TIMESTAMP
                WHERE content_hash = %s
                """,
                (text, object_name, file_type, content_hash),
            )
            logger.info(f"[ppt] Updated existing record in unstructured_documents for {object_name}")
        else:
            # Insert new record
            cursor.execute(
                """
                INSERT INTO unstructured_documents 
                    (object_name, file_type, text_content, content_hash)
                VALUES (%s, %s, %s, %s)
                """,
                (object_name, file_type, text, content_hash),
            )
            logger.info(f"[ppt] Inserted new record into unstructured_documents for {object_name}")
        
        pg_conn.commit()
        
    except Exception as e:
        pg_conn.rollback()
        logger.exception(f"[ppt] Failed saving to unstructured_documents: {e}")
        logger.warning("[ppt] Continuing pipeline despite unstructured_documents save failure")
    finally:
        cursor.close()


def extract_text_from_slide(slide) -> Dict[str, Any]:
    """
    Extract all text content from a single slide.
    Returns structured data including title, body text, and notes.
    """
    slide_data = {
        'title': '',
        'body': [],
        'notes': '',
        'tables': []
    }
    
    try:
        # Extract shapes (text boxes, titles, etc.)
        for shape in slide.shapes:
            try:
                # Handle text frames
                if hasattr(shape, 'text') and shape.text:
                    text = shape.text.strip()
                    
                    # Try to identify if it's a title
                    if hasattr(shape, 'name') and 'title' in shape.name.lower():
                        slide_data['title'] = text
                    else:
                        slide_data['body'].append(text)
                
                # Handle tables
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_data = extract_table_from_shape(shape)
                    if table_data is not None:
                        slide_data['tables'].append(table_data)
                        
            except Exception as e:
                logger.warning(f"[ppt] Error extracting from shape: {e}")
                continue
        
        # Extract speaker notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    slide_data['notes'] = notes_text.strip()
            except Exception as e:
                logger.warning(f"[ppt] Error extracting notes: {e}")
                
    except Exception as e:
        logger.warning(f"[ppt] Error processing slide: {e}")
    
    return slide_data


def extract_table_from_shape(shape) -> Optional[pd.DataFrame]:
    """Extract table data from a table shape and convert to DataFrame."""
    try:
        if not hasattr(shape, 'table'):
            return None
        
        table = shape.table
        rows = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                try:
                    cell_text = cell.text.strip() if cell.text else ''
                    row_data.append(cell_text)
                except Exception as e:
                    logger.warning(f"[ppt] Error reading cell: {e}")
                    row_data.append('')
            rows.append(row_data)
        
        if not rows or len(rows) < 2:
            return None
        
        # Use first row as headers
        df = pd.DataFrame(rows[1:], columns=rows[0])
        return df
        
    except Exception as e:
        logger.warning(f"[ppt] Error extracting table: {e}")
        return None


def extract_images_from_pptx(data: bytes, file_root: str) -> List[Dict[str, Any]]:
    """
    Extract images from PowerPoint file.
    Returns list of image metadata (could be used to save images to MinIO).
    """
    images = []
    
    try:
        # PowerPoint files are ZIP archives
        with ZipFile(io.BytesIO(data)) as zip_file:
            # Images are typically in ppt/media/ folder
            for file_info in zip_file.filelist:
                if file_info.filename.startswith('ppt/media/'):
                    try:
                        image_data = zip_file.read(file_info.filename)
                        image_name = file_info.filename.split('/')[-1]
                        
                        images.append({
                            'name': image_name,
                            'size': len(image_data),
                            'data': image_data,
                            'path': f"processed/unstructured/ppt-images/{file_root}/{image_name}"
                        })
                    except Exception as e:
                        logger.warning(f"[ppt] Error extracting image {file_info.filename}: {e}")
                        
    except Exception as e:
        logger.warning(f"[ppt] Error extracting images from ZIP: {e}")
    
    return images


def _process_extracted_table(
    minio_client,
    bucket_name: str,
    table_df: pd.DataFrame,
    table_key: str,
    pg_conn,
    catalog_updater,
):
    """
    Process a single extracted table from PowerPoint.
    Similar to PDF table processing.
    """
    try:
        # Import from structured_pipeline for consistency
        from pipelines.structured_pipeline import (
            sanitize_table_name,
            sanitize_column_name,
            infer_postgres_type,
            normalize_dataframe
        )
        
        # Normalize the dataframe
        table_df = normalize_dataframe(table_df)
        
        # Remove empty rows/columns
        table_df = table_df.dropna(how='all', axis=0)
        table_df = table_df.dropna(how='all', axis=1)
        table_df = table_df.replace('', None)
        
        if table_df.empty or len(table_df.columns) == 0:
            logger.warning(f"[ppt] Table is empty after normalization, skipping: {table_key}")
            return
        
        # Upload CSV to MinIO
        csv_bytes = table_df.to_csv(index=False).encode("utf-8")
        minio_client.put_object(
            bucket_name,
            table_key,
            io.BytesIO(csv_bytes),
            length=len(csv_bytes),
            content_type="text/csv",
        )
        logger.info(f"[ppt] Uploaded table CSV -> {table_key}")

        # Prepare for PostgreSQL
        table_name = sanitize_table_name(table_key)
        
        # Sanitize column names
        table_df.columns = [sanitize_column_name(col) for col in table_df.columns]
        
        # Handle duplicate column names
        cols = pd.Series(table_df.columns)
        for dup in cols[cols.duplicated()].unique():
            cols[cols == dup] = [f"{dup}_{i}" if i != 0 else dup for i in range(sum(cols == dup))]
        table_df.columns = cols
        
        # Load to PostgreSQL
        cursor = pg_conn.cursor()
        
        try:
            # Drop existing table
            cursor.execute(f'DROP TABLE IF EXISTS "{table_name}"')
            
            # Create table with inferred types
            column_defs = []
            for col in table_df.columns:
                try:
                    pg_type = infer_postgres_type(table_df[col])
                except Exception as type_error:
                    logger.warning(f"[ppt] Failed to infer type for column {col}: {type_error}. Using TEXT.")
                    pg_type = 'TEXT'
                column_defs.append(f'"{col}" {pg_type}')
            
            create_sql = f'CREATE TABLE "{table_name}" ({", ".join(column_defs)})'
            cursor.execute(create_sql)
            logger.info(f"[ppt] Created table: {table_name}")
            
            # Bulk insert
            buffer = io.StringIO()
            table_df.to_csv(buffer, index=False, header=False, sep='\t', na_rep='\\N')
            buffer.seek(0)
            
            cursor.copy_expert(
                f'COPY "{table_name}" FROM STDIN WITH (FORMAT CSV, DELIMITER E\'\\t\', NULL \'\\N\')',
                buffer
            )
            
            pg_conn.commit()
            logger.info(f"[ppt] Loaded {len(table_df)} rows into {table_name}")
            
        except Exception as db_error:
            pg_conn.rollback()
            logger.exception(f"[ppt] Database operation failed for {table_name}: {db_error}")
            raise
        finally:
            cursor.close()

        # Update catalog
        try:
            catalog_updater(
                object_name=table_key,
                object_size=len(csv_bytes),
                file_format='csv',
                row_count=len(table_df),
                text_extracted=False,
                content_hash=None
            )
        except Exception as catalog_error:
            logger.warning(f"[ppt] Failed to update catalog for {table_key}: {catalog_error}")

    except Exception as e:
        logger.exception(f"[ppt] Failed processing table {table_key}: {e}")
        if pg_conn:
            pg_conn.rollback()


def process_minio_object(
    minio_client: Minio,
    bucket_name: str,
    object_name: str,
    pg_conn,
    catalog_updater,
):
    """
    Main entry point for PowerPoint processing.

    Steps:
    1. Downloads raw PPTX from MinIO
    2. Computes SHA256 hash and checks for duplicates
    3. Extracts text from all slides (titles, body, notes)
    4. Saves text to:
       - MinIO: processed/unstructured/text-extracted/<file>.txt
       - MinIO: processed/unstructured/ppt-structured/<file>.json (structured data)
       - Postgres: unstructured_documents table
    5. Extracts tables from slides:
       - Uploads each as CSV to processed/structured/ppt-tables/<file>_slide_X_table_Y.csv
       - Creates PostgreSQL table for each
       - Updates catalog for each table
    6. Optionally extracts images and saves to MinIO
    7. Updates catalog for original PPTX file
    """
    logger.info(f"[ppt] Processing {object_name}")
    
    response = None
    data = None
    
    try:
        # 1. Download from MinIO
        response = minio_client.get_object(bucket_name, object_name)
        data = response.read()
        
    except Exception as e:
        logger.exception(f"[ppt] Failed to download {object_name} from MinIO: {e}")
        raise
    finally:
        if response:
            response.close()
            response.release_conn()

    if not data:
        logger.error(f"[ppt] No data downloaded for {object_name}")
        return

    file_size = len(data)
    file_hash = calculate_file_hash(data)
    file_root = object_name.split("/")[-1].rsplit(".", 1)[0]

    # 2. Check for duplicates
    if is_duplicate(pg_conn, file_hash):
        logger.info(
            f"[ppt] Duplicate detected via content_hash, "
            f"skipping heavy processing: {object_name}"
        )
        try:
            catalog_updater(
                object_name=object_name,
                object_size=file_size,
                file_format="pptx",
                row_count=0,
                text_extracted=False,
                content_hash=file_hash,
            )
        except Exception as e:
            logger.warning(f"[ppt] Failed catalog update for duplicate: {e}")
        return

    # 3. Extract content from PowerPoint
    logger.info(f"[ppt] Extracting content from PowerPoint...")
    
    slides_data = []
    full_text_parts = []
    table_count = 0
    
    try:
        prs = Presentation(io.BytesIO(data))
        slide_count = len(prs.slides)
        
        logger.info(f"[ppt] Processing {slide_count} slides...")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            try:
                # Extract all content from this slide
                slide_data = extract_text_from_slide(slide)
                slide_data['slide_number'] = slide_idx
                slides_data.append(slide_data)
                
                # Build full text for search
                text_parts = []
                if slide_data['title']:
                    text_parts.append(f"SLIDE {slide_idx}: {slide_data['title']}")
                if slide_data['body']:
                    text_parts.extend(slide_data['body'])
                if slide_data['notes']:
                    text_parts.append(f"Notes: {slide_data['notes']}")
                
                full_text_parts.extend(text_parts)
                
                # Process tables from this slide
                for table_idx, table_df in enumerate(slide_data['tables'], 1):
                    if table_df is not None and not table_df.empty:
                        table_key = (
                            f"processed/structured/ppt-tables/"
                            f"{file_root}_slide{slide_idx}_table{table_idx}.csv"
                        )
                        
                        _process_extracted_table(
                            minio_client=minio_client,
                            bucket_name=bucket_name,
                            table_df=table_df,
                            table_key=table_key,
                            pg_conn=pg_conn,
                            catalog_updater=catalog_updater,
                        )
                        
                        table_count += 1
                        logger.info(
                            f"[ppt] Processed table {table_count}: "
                            f"slide {slide_idx}, {len(table_df)} rows, {len(table_df.columns)} columns"
                        )
                
                if slide_idx % 10 == 0:
                    logger.info(f"[ppt] Processed {slide_idx}/{slide_count} slides")
                    
            except Exception as slide_error:
                logger.warning(f"[ppt] Error processing slide {slide_idx}: {slide_error}")
        
        logger.info(f"[ppt] Extracted content from {slide_count} slides, {table_count} tables")
        
    except Exception as e:
        logger.exception(f"[ppt] Failed to extract content from PowerPoint: {e}")
        return

    # 4. Save full text to MinIO
    full_text = "\n\n".join(full_text_parts)
    
    if full_text.strip():
        try:
            text_bytes = full_text.encode("utf-8")
            text_path = f"processed/unstructured/text-extracted/{file_root}.txt"

            minio_client.put_object(
                bucket_name,
                text_path,
                io.BytesIO(text_bytes),
                length=len(text_bytes),
                content_type="text/plain",
            )
            logger.info(f"[ppt] Uploaded extracted text -> {text_path}")

        except Exception as e:
            logger.exception(f"[ppt] Failed to upload text to MinIO: {e}")

        # Save to database
        if pg_conn is not None:
            try:
                _save_unstructured_doc(
                    pg_conn=pg_conn,
                    object_name=object_name,
                    file_type="pptx",
                    text=full_text,
                    content_hash=file_hash,
                )
            except Exception as e:
                logger.warning(f"[ppt] Failed to save unstructured doc: {e}")
    
    # 5. Save structured JSON data
    try:
        structured_json = json.dumps(slides_data, indent=2, ensure_ascii=False)
        json_bytes = structured_json.encode("utf-8")
        json_path = f"processed/unstructured/ppt-structured/{file_root}.json"
        
        minio_client.put_object(
            bucket_name,
            json_path,
            io.BytesIO(json_bytes),
            length=len(json_bytes),
            content_type="application/json",
        )
        logger.info(f"[ppt] Uploaded structured JSON -> {json_path}")
        
    except Exception as e:
        logger.warning(f"[ppt] Failed to save structured JSON: {e}")

    # 6. Extract and save images (optional)
    try:
        images = extract_images_from_pptx(data, file_root)
        
        for img in images:
            try:
                minio_client.put_object(
                    bucket_name,
                    img['path'],
                    io.BytesIO(img['data']),
                    length=img['size'],
                    content_type='application/octet-stream',
                )
                logger.info(f"[ppt] Saved image: {img['name']}")
            except Exception as img_error:
                logger.warning(f"[ppt] Failed to save image {img['name']}: {img_error}")
        
        if images:
            logger.info(f"[ppt] Extracted and saved {len(images)} images")
            
    except Exception as e:
        logger.warning(f"[ppt] Failed to extract images: {e}")

    # 7. Update catalog for original PPTX file
    try:
        prs = Presentation(io.BytesIO(data))
        props = prs.core_properties
        ppt_metadata = {
            'author': props.author,
            'created': props.created.isoformat() if props.created else None,
            'modified': props.modified.isoformat() if props.modified else None,
            'title': props.title,
            'subject': props.subject,
            'category': props.category,
            'comments': props.comments,
            'slide_count': len(prs.slides),
            'table_count': table_count,
            'image_count': len(images)
        }

        catalog_updater(
            object_name=object_name,
            object_size=file_size,
            file_format="pptx",
            row_count=table_count,
            text_extracted=bool(full_text.strip()),
            content_hash=file_hash,
            metadata=ppt_metadata  # NEW
        )
        logger.info(f"[ppt] Updated catalog for {object_name}")
    except Exception as e:
        logger.exception(f"[ppt] Failed catalog update for {object_name}: {e}")

    logger.info(
        f"[ppt] ✅ Completed processing: {object_name} | "
        f"slides={len(slides_data)} | tables={table_count} | "
        f"text_chars={len(full_text)}"
    )